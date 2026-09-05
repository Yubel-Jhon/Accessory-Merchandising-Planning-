#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""企划导出（v0.3）：把「企划盘」（多 SKU）渲染成完整 deck——HTML 长页 + 多 slide PPTX。

v0.1 是 1 页拼图；v0.2 对齐 COLE HAAN FW26 结构升级为 8 页单款 deck；
v0.3 升级为多 SKU 企划盘（PRD v2 二期 / 优化报告 P1-A）：

  P01 封面        方向 × N 款 × 「AI 企划」副标题
  P02 企划方法    4 步流程 + 工具/耗时（静态模板）
  P03 人群画像    3 类人群（direction.personas）
  P04 产品结构总表 全部款：品类 / 成分 / 规格 / 双价格 / 对标（企划盘聚合）
  P05+ 逐款页×N   图组 + 参数条（成分/规格/细度/工艺/功能）+ 双价格 + 设计方向
  P06+ 演变对比   before/after（该款有演变记录时，紧跟其逐款页）
  P07+ 品类矩阵   全部款拼版 + 会员价（一页看全盘）
  P08+ AI 出图体系 行=款 × 列=图类型（每款已生成的图）
  P09 开发日历    8 节点时间线 + AI 压缩标注（静态模板）
  P10 尾页        工具链 + 真实耗时拆解（plan.timing）+ THANKS

plan 结构（v0.3；兼容 v0.1/v0.2 单款字段，缺失逐项回退）：
{
  "direction": str, "retailer": str,
  "skus": [                                        # 企划盘：多款
    { "sku": {全字段 v2}, "color": str,
      "selected": {image_type: 图片绝对路径},
      "variation": {before, after, axis, change} | null },
  ],
  # —— v0.1/v0.2 单款兼容（skus 缺失时由此合成一条目）——
  "sku": {...}, "product": {...}, "selected": {...}, "variation": {...},
  "personas": [{name, desc, need}, ...],
  "selection_logic": {标题: 正文},
  "timing": {evolve_sec, images_sec, total_min} | {}
}
"""
import base64
import os
from datetime import date

TYPE_LABEL = {
    "white_bg": "白底图 · 产品身份证",
    "studio": "商拍图 · 质感定调",
    "lifestyle": "氛围图 · 场景代入",
    "detail": "细节图 · 品质论证",
    "fabric": "面料图 · 材料透明",
    "model": "模特图 · 上身效果",
}
LAYOUT_ORDER = ["white_bg", "detail", "fabric", "studio", "lifestyle", "model"]

AXIS_LABEL = {"color": "改色", "detail": "改细节", "silhouette": "改廓形"}

# FW26/27 开发日历（静态模板，8 节点，对齐行业开发节奏）
CALENDAR = [
    ("3/15-3/31", "企划"), ("4/1-4/15", "设计企划"), ("4/16-5/15", "产品开发"),
    ("5/16-5/20", "定版"), ("6/1-6/10", "下单"), ("6/11-7/20", "备料·打样·测试·生产"),
    ("7/21-7/25", "验货·出货"), ("8/1-8/10", "进店销售"),
]

# 企划方法 4 步（静态模板）
METHOD_STEPS = [
    ("01", "锚点选款", "拆解畅销款三要素：原料等级 / 无 logo / 价格带", "qwen-vl 识别 + 对标库"),
    ("02", "单轴演变", "改色 / 改细节 / 改廓形，锁 DNA 只动一个轴", "qwen-image 图生图 · 约 1–3 分钟/张"),
    ("03", "AI 出图", "白底/商拍/氛围/细节/面料/模特，一套 SKU 全类图覆盖", "qwen-image 图生图"),
    ("04", "企划组装", "参数条 + 双价格 + 开发日历 → 买手推介 deck", "exporter 自动编排"),
]

TOOLCHAIN = "图生图/文生图 qwen-image-3.0-pro · 锚点识别 qwen-vl-max"

FALLBACK_PERSONAS = [
    {"name": "城市老钱", "desc": "沉稳、内敛、非凡成就", "need": "看得见的原料等级，看不见的 logo"},
    {"name": "都市菁英", "desc": "质感、精致、全球视野", "need": "通勤场景的体面，价格不用解释"},
    {"name": "Z世代新贵", "desc": "自信、独立、充满激情", "need": "大牌同源质感，能晒的静奢感"},
]


# ---------- 数据整理 ----------

def _normalize_entries(plan):
    """企划盘（v0.3 多款）→ 条目列表；v0.1/v0.2 单款 plan 由旧字段合成一条目。空款剔除。"""
    entries = []
    for e in plan.get("skus") or []:
        if e.get("selected"):
            entries.append({
                "sku": e.get("sku") or {},
                "color": e.get("color", ""),
                "selected": e["selected"],
                "variation": e.get("variation"),
            })
    if not entries and plan.get("selected"):  # 旧单款兼容
        p = plan.get("product") or {}
        entries.append({
            "sku": plan.get("sku") or {},
            "color": p.get("color", ""),
            "selected": plan["selected"],
            "variation": plan.get("variation"),
        })
    return entries


def _entry_view(entry):
    """单款条目 → 取值视图（v2 sku 字段 + v0.1 product 字段全兼容）。"""
    sku = entry.get("sku") or {}
    price = sku.get("price") or {}
    msrp, ws, cur = price.get("msrp"), price.get("wholesale"), price.get("currency", "¥")
    return {
        "category": sku.get("name") or "商品",
        "benchmark": sku.get("benchmark") or "",
        "material": sku.get("material") or "",
        "craft": sku.get("craft") or "",
        "composition": sku.get("composition") or sku.get("material") or "",
        "spec": sku.get("spec") or "待补充",
        "fineness": sku.get("fineness") or "",
        "attributes": sku.get("attributes") or "",
        "design_directions": sku.get("design_directions") or [],
        "color": entry.get("color", ""),
        "msrp_str": f'{cur}{msrp}' if msrp else "待校准",
        "wholesale_str": f'{cur}{ws}' if ws else "待校准",
        "price_str": (f'会员价 {cur}{msrp} ｜ 供货价 {cur}{ws}' if msrp and ws else "价格待校准"),
    }


def _timing_str(timing):
    if not timing:
        return "由「商品企划 Agent」全流程编排生成"
    parts = []
    if timing.get("evolve_sec"):
        parts.append(f'单轴演变 {timing["evolve_sec"]} 秒')
    if timing.get("images_sec"):
        parts.append(f'出图 {timing["images_sec"]} 秒')
    if timing.get("total_min"):
        parts.append(f'合计约 {timing["total_min"]} 分钟')
    return "Agent 实测耗时：" + " · ".join(parts) if parts else "由「商品企划 Agent」全流程编排生成"


def _design_dirs_str(view):
    if not view["design_directions"]:
        return ""
    return " ｜ ".join(f'{i+1}. {AXIS_LABEL.get(d["axis"], d["axis"])}：{d["desc"]}'
                       for i, d in enumerate(view["design_directions"]))


def _b64(path):
    with open(path, "rb") as f:
        return "data:image/jpeg;base64," + base64.b64encode(f.read()).decode()


# ---------- HTML 导出 ----------

def render_html(plan):
    entries = _normalize_entries(plan)
    if not entries:
        raise ValueError("plan 里没有任何已选图片（skus/selected 均为空）")
    views = [_entry_view(e) for e in entries]
    personas = plan.get("personas") or FALLBACK_PERSONAS

    chips = [
        ("方向", plan.get("direction", "")), ("款式数", f"{len(entries)}"),
        ("零售商", plan.get("retailer", "")), ("整体风格", plan.get("plan_style", "")),
        ("日期", date.today().isoformat()),
    ]
    chip_html = "".join(f'<span class="chip"><b>{k}</b>{v}</span>' for k, v in chips if v)

    steps_html = "".join(
        f'<div class="card"><div class="stepno">{no}</div><h3>{name}</h3><p>{desc}</p>'
        f'<p class="tool">🛠 {tool}</p></div>'
        for no, name, desc, tool in METHOD_STEPS)

    personas_html = "".join(
        f'<div class="card"><h3>{p["name"]}</h3><p>{p["desc"]}</p><p class="tool">要什么：{p.get("need","")}</p></div>'
        for p in personas)

    # P04 总表
    table_rows = "".join(
        f'<tr><td class="tname">{v["category"]}</td><td>{v["composition"]}</td><td>{v["spec"]}</td>'
        f'<td>{v["msrp_str"]}</td><td>{v["wholesale_str"]}</td><td>{v["benchmark"]}</td></tr>'
        for v in views)

    # P05+ 逐款页 + 演变
    entries_html = ""
    for i, (e, v) in enumerate(zip(entries, views)):
        layout = [t for t in LAYOUT_ORDER if t in e["selected"]]
        imgs = {t: _b64(e["selected"][t]) for t in layout}
        hero_key = "lifestyle" if "lifestyle" in e["selected"] else (layout[0] if layout else None)
        hero = imgs.get(hero_key, "")
        sec = [t for t in layout if t != hero_key][:2]
        sec_html = "".join(f'<figure class="gitem"><img src="{imgs[t]}"/>'
                           f'<figcaption>{TYPE_LABEL[t].split(" ·")[0]}</figcaption></figure>' for t in sec)
        params = [("成分", v["composition"]), ("规格", v["spec"]), ("细度", v["fineness"]),
                  ("工艺", v["craft"]), ("功能", v["attributes"]), ("对标", v["benchmark"])]
        param_html = "".join(f'<div class="prow"><b>{k}</b><span>{val}</span></div>' for k, val in params if val)
        dd = _design_dirs_str(v)
        entries_html += (
            f'<h2 class="sec">逐款企划 {i+1}/{len(entries)} · {v["category"]}'
            + (f'（{v["color"]}）' if v["color"] else '') + '</h2>'
            + f'<div class="entry"><div class="entry-imgs">'
              f'<img class="hero" src="{hero}">{sec_html}</div>'
              f'<div class="ppanel">{param_html}'
              f'<div class="price"><div><div class="pv">{v["msrp_str"]}</div><div class="lab">MSRP · 会员价</div></div>'
              f'<div><div class="pv">{v["wholesale_str"]}</div><div class="lab">WHOLESALE · 供货价</div></div></div>'
              + (f'<div class="dd"><b>设计方向：</b>{dd}</div>' if dd else '')
              + '</div></div>')
        ve = e.get("variation")
        if ve:
            axis_label = AXIS_LABEL.get(ve.get("axis"), ve.get("axis", ""))
            entries_html += (
                '<div class="compare">'
                f'<figure><img src="{_b64(ve["before"])}"/><figcaption>畅销款（起点）</figcaption></figure>'
                '<span class="arrow">→</span>'
                f'<figure><img src="{_b64(ve["after"])}"/>'
                f'<figcaption>演变款（{axis_label}：{ve.get("change","")}）</figcaption></figure>'
                '</div>')

    # 品类矩阵：一页看全盘
    matrix_html = "".join(
        f'<figure class="mcell"><img src="{_b64(next(iter(e["selected"].values())))}"/>'
        f'<figcaption><b>{v["category"]}</b><span>{v["msrp_str"]}</span></figcaption></figure>'
        for e, v in zip(entries, views))

    # AI 出图体系：行=款 × 列=图类型
    sys_rows = ""
    for e, v in zip(entries, views):
        layout = [t for t in LAYOUT_ORDER if t in e["selected"]]
        row_imgs = "".join(f'<img src="{_b64(e["selected"][t])}" title="{TYPE_LABEL[t]}">' for t in layout)
        sys_rows += f'<div class="sysrow"><div class="sysname">{v["category"]}</div><div class="sysimgs">{row_imgs}</div></div>'

    cal_html = "".join(f'<div class="cal"><b>{d}</b><span>{name}</span></div>' for d, name in CALENDAR)
    cover = plan.get("cover")
    cover_ok = cover and os.path.isfile(cover)
    # P01 hero：有封面图（用户从氛围图挑的/再生成的）用它；没有回退到首款首图
    hero0 = _b64(cover) if cover_ok else _b64(next(iter(entries[0]["selected"].values())))

    return f"""<!DOCTYPE html>
<html lang="zh-CN"><head><meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{plan.get("direction","")} · 产品企划盘</title><style>
:root{{--ink:#2C2C2C;--brown:#8B7355;--paper:#F5F1EA;--card:#FFF;--line:#E5DED2}}
*{{box-sizing:border-box;margin:0;padding:0}}
body{{font-family:"Segoe UI","PingFang SC","Microsoft YaHei",sans-serif;background:var(--paper);color:var(--ink);line-height:1.65}}
.wrap{{max-width:1120px;margin:0 auto;padding:40px 28px 80px}}
header{{text-align:center;padding:28px 0 8px}}
header .eyebrow{{letter-spacing:.35em;font-size:12px;color:var(--brown);text-transform:uppercase;margin-bottom:10px}}
header h1{{font-family:Georgia,"Songti SC",serif;font-weight:500;font-size:40px}}
header .sub{{color:#7a6f5f;margin-top:10px;font-size:15px}}
.badge{{display:inline-block;margin-top:14px;padding:6px 18px;border:1px solid var(--brown);border-radius:999px;color:var(--brown);font-size:13px}}
.chips{{display:flex;flex-wrap:wrap;gap:10px;justify-content:center;margin:26px 0 6px}}
.chip{{background:var(--card);border:1px solid var(--line);border-radius:999px;padding:6px 14px;font-size:13px;color:#4a4238}}
.chip b{{color:var(--brown);margin-right:6px;font-weight:600}}
.hero{{margin:34px 0 26px;border-radius:6px;overflow:hidden;box-shadow:0 12px 40px rgba(44,44,44,.14)}}
.hero img{{width:100%;display:block}}
h2.sec{{font-family:Georgia,"Songti SC",serif;font-weight:500;font-size:22px;margin:46px 0 18px;display:flex;align-items:center;gap:12px}}
h2.sec::before{{content:"";width:26px;height:2px;background:var(--brown);display:inline-block}}
.cards{{display:grid;grid-template-columns:repeat(auto-fit,minmax(220px,1fr));gap:14px}}
.card{{background:var(--card);border:1px solid var(--line);border-radius:6px;padding:18px 20px}}
.card h3{{font-size:14px;color:var(--brown);letter-spacing:.08em;margin-bottom:8px;font-weight:600}}
.card p{{font-size:14px;color:#3f3930}}
.card .stepno{{font-family:Georgia,serif;font-size:26px;color:var(--brown);margin-bottom:6px}}
.card .tool{{margin-top:8px;font-size:12px;color:#8a7f6f}}
table.skutable{{width:100%;border-collapse:collapse;background:var(--card);border:1px solid var(--line);border-radius:6px;overflow:hidden;font-size:13px}}
.skutable th{{background:#f0e9dc;color:var(--brown);padding:10px 12px;text-align:left;font-weight:600;font-size:12px;letter-spacing:.05em}}
.skutable td{{padding:10px 12px;border-top:1px solid var(--line)}}
.skutable .tname{{font-weight:600;color:var(--ink)}}
.entry{{display:grid;grid-template-columns:1.05fr 1fr;gap:16px;align-items:start;margin-bottom:8px}}
.entry-imgs{{display:grid;grid-template-columns:1fr 1fr;gap:10px}}
.entry-imgs .hero{{grid-column:1/-1;width:100%;aspect-ratio:1/1;object-fit:cover;border-radius:6px;border:1px solid var(--line)}}
.gitem{{background:var(--card);border:1px solid var(--line);border-radius:6px;overflow:hidden}}
.gitem img{{width:100%;aspect-ratio:1/1;object-fit:cover;display:block}}
.gitem figcaption{{padding:8px 10px;font-size:12px;color:#5b5145;border-top:1px solid var(--line);background:#fbf9f5}}
.ppanel{{background:var(--card);border:1px solid var(--line);border-radius:6px;padding:20px}}
.prow{{display:flex;gap:12px;padding:7px 0;border-bottom:1px dashed var(--line);font-size:14px}}
.prow b{{color:var(--brown);min-width:3em;font-weight:600}}
.price{{display:flex;gap:24px;margin:14px 0 6px}}
.price .pv{{font-family:Georgia,serif;font-size:26px;font-weight:600}}
.price .lab{{font-size:12px;color:#8a7f6f}}
.dd{{margin-top:10px;font-size:13px;color:#5b5145}}
.matrix{{display:grid;grid-template-columns:repeat(auto-fit,minmax(200px,1fr));gap:14px}}
.mcell{{background:var(--card);border:1px solid var(--line);border-radius:6px;overflow:hidden}}
.mcell img{{width:100%;aspect-ratio:1/1;object-fit:cover;display:block}}
.mcell figcaption{{display:flex;justify-content:space-between;align-items:center;padding:10px 12px;font-size:13px;border-top:1px solid var(--line)}}
.mcell figcaption span{{font-family:Georgia,serif;font-weight:600;color:var(--brown)}}
.sysrow{{display:grid;grid-template-columns:140px 1fr;gap:12px;align-items:center;background:var(--card);border:1px solid var(--line);border-radius:6px;padding:10px;margin-bottom:10px}}
.sysname{{font-size:13px;font-weight:600;color:var(--brown)}}
.sysimgs{{display:flex;gap:8px;overflow:hidden}}
.sysimgs img{{width:86px;height:86px;object-fit:cover;border-radius:4px;border:1px solid var(--line)}}
.timeline{{display:grid;grid-template-columns:repeat(8,1fr);gap:8px}}
.cal{{background:var(--card);border:1px solid var(--line);border-radius:6px;padding:10px;text-align:center}}
.cal b{{display:block;font-size:12px;color:var(--brown)}}
.cal span{{font-size:12px}}
.airule{{margin-top:16px;padding:12px 16px;border-left:3px solid var(--brown);background:var(--card);font-size:14px}}
.compare{{display:grid;grid-template-columns:1fr auto 1fr;gap:14px;align-items:center;margin:14px 0 26px;max-width:640px}}
.compare figure{{background:var(--card);border:1px solid var(--line);border-radius:6px;overflow:hidden}}
.compare img{{width:100%;aspect-ratio:1/1;object-fit:cover;display:block}}
.compare figcaption{{padding:9px 12px;font-size:13px;color:#5b5145;border-top:1px solid var(--line)}}
.compare .arrow{{font-size:30px;color:var(--brown)}}
footer{{margin-top:60px;padding-top:20px;border-top:1px solid var(--line);font-size:12px;color:#a49a8b;text-align:center}}
@media(max-width:800px){{.entry{{grid-template-columns:1fr}}}}
</style></head><body><div class="wrap">
<header><div class="eyebrow">Product Planning · {plan.get("retailer","")}</div>
<h1>{plan.get("direction","")} · 产品企划盘</h1>
<div class="sub">{len(entries)} 个款式 ｜ {plan.get("retailer","")} 渠道 ｜ {date.today().isoformat()}</div>
<div class="badge">⚡ AI 企划 · 全盘概念图 45 分钟</div></header>
<div class="chips">{chip_html}</div>
<div class="hero"><img src="{hero0}"></div>
<h2 class="sec">企划方法：数据选款 → AI 演变 → 买手推介</h2><div class="cards">{steps_html}</div>
<h2 class="sec">人群画像</h2><div class="cards">{personas_html}</div>
<h2 class="sec">产品结构总表</h2>
<table class="skutable"><tr><th>品类</th><th>成分</th><th>规格</th><th>会员价</th><th>供货价</th><th>对标</th></tr>
{table_rows}</table>
<h2 class="sec">逐款企划</h2>
{entries_html}
<h2 class="sec">品类矩阵 · 一页看全盘</h2><div class="matrix">{matrix_html}</div>
<h2 class="sec">AI 出图体系 · 行=款 × 列=图类型</h2>{sys_rows}
<h2 class="sec">开发日历</h2><div class="timeline">{cal_html}</div>
<div class="airule">⚡ AI 压缩：设计出图环节由「人天级」→「分钟级」，今天这套 deck 即 Agent 实跑产物。</div>
<footer>{TOOLCHAIN}<br>{_timing_str(plan.get("timing"))} ｜ {date.today().isoformat()}</footer>
</div></body></html>"""


# ---------- PPTX 导出（多 slide 企划盘 deck） ----------

def render_pptx(plan, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from PIL import Image

    INK = RGBColor(0x2C, 0x2C, 0x2C)
    BROWN = RGBColor(0x8B, 0x73, 0x55)
    GREY = RGBColor(0x5B, 0x51, 0x45)
    LINE = RGBColor(0xE5, 0xDE, 0xD2)
    PAPER = RGBColor(0xF5, 0xF1, 0xEA)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    entries = _normalize_entries(plan)
    if not entries:
        raise ValueError("plan 里没有任何已选图片（skus/selected 均为空）")
    views = [_entry_view(e) for e in entries]
    personas = plan.get("personas") or FALLBACK_PERSONAS

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def box(x, y, w, h, fill=None, line=None, slide=None):
        slide = slide or prs.slides[-1]
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid() if fill else shp.fill.background()
        if fill:
            shp.fill.fore_color.rgb = fill
        if line:
            shp.line.color.rgb = line
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def text(x, y, w, h, runs, align=None, slide=None):
        slide = slide or prs.slides[-1]
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, (txt, size, color, bold) in enumerate(runs):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            r = para.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = "Microsoft YaHei"
        return tb

    def new_slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box(0, 0, 13.333, 7.5, fill=PAPER, slide=s)
        box(0, 0, 13.333, 0.12, fill=BROWN, slide=s)
        return s

    def page_title(s, title, sub=""):
        text(0.55, 0.30, 12.2, 0.8, [(title, 28, INK, True)], slide=s)
        if sub:
            text(0.55, 1.0, 12.2, 0.4, [(sub, 13, BROWN, False)], slide=s)

    def fit_img(path, cx, cy, cw, ch, slide=None):
        with Image.open(path) as im:
            w, h = im.size
        ar = w / h; box_ar = cw / ch
        if ar > box_ar:
            dw, dh = cw, cw / ar
        else:
            dh, dw = ch, ch * ar
        dx = cx + (cw - dw) / 2; dy = cy + (ch - dh) / 2
        (slide or prs.slides[-1]).shapes.add_picture(path, Inches(dx), Inches(dy), Inches(dw), Inches(dh))

    def card(x, y, w, h, slide=None):
        return box(x, y, w, h, fill=WHITE, line=LINE, slide=slide)

    # ---- P01 封面 ----
    s = new_slide()
    cover = plan.get("cover")
    cover_ok = cover and os.path.isfile(cover)
    plan_style = plan.get("plan_style", "")
    if cover_ok:
        # 有封面图：左文右图版式
        text(0.9, 2.1, 6.0, 0.5, [("PRODUCT PLANNING ·  " + plan.get("retailer", ""), 14, BROWN, False)], slide=s)
        text(0.9, 2.6, 6.0, 1.2, [(f'{plan.get("direction","")} · 产品企划盘', 40, INK, True)], slide=s)
        text(0.9, 3.9, 6.0, 0.5,
             [(f'{len(entries)} 个款式 ｜ {plan.get("retailer","")} 渠道 ｜ ' + date.today().isoformat(), 14, GREY, False)], slide=s)
        if plan_style:
            text(0.9, 4.4, 6.0, 0.4, [("整体企划风格：" + plan_style, 12, BROWN, False)], slide=s)
        card(0.9, 5.0, 3.6, 0.55, slide=s)
        text(0.9, 5.08, 3.6, 0.4, [("⚡ AI 企划 · 全盘概念图 45 分钟", 14, BROWN, True)], slide=s)
        fit_img(cover, 7.0, 1.9, 5.5, 3.7, slide=s)
        text(7.0, 5.68, 5.5, 0.35, [("deck 封面 · AI 生成", 10, GREY, False)], slide=s)
    else:
        text(0.9, 2.1, 11.5, 0.5, [("PRODUCT PLANNING ·  " + plan.get("retailer", ""), 14, BROWN, False)], slide=s)
        text(0.9, 2.6, 11.5, 1.2, [(f'{plan.get("direction","")} · 产品企划盘', 44, INK, True)], slide=s)
        text(0.9, 3.9, 11.5, 0.5,
             [(f'{len(entries)} 个款式 ｜ {plan.get("retailer","")} 渠道 ｜ ' + date.today().isoformat(), 16, GREY, False)], slide=s)
        if plan_style:
            text(0.9, 4.35, 11.5, 0.3, [("整体企划风格：" + plan_style, 12, BROWN, False)], slide=s)
        card(0.9, 4.7, 3.6, 0.55, slide=s)
        text(0.9, 4.78, 3.6, 0.4, [("⚡ AI 企划 · 全盘概念图 45 分钟", 14, BROWN, True)], slide=s)
    text(0.9, 6.7, 11.5, 0.4, [(date.today().isoformat() + " ｜ 商品企划 Agent 生成", 11, GREY, False)], slide=s)

    # ---- P02 企划方法 ----
    s = new_slide()
    page_title(s, "企划方法：数据选款 → AI 演变 → 买手推介")
    cw, ch_, gap = 2.95, 4.6, 0.15
    for i, (no, name, desc, tool) in enumerate(METHOD_STEPS):
        x = 0.55 + i * (cw + gap)
        card(x, 1.7, cw, ch_, slide=s)
        text(x + 0.2, 1.95, cw - 0.4, 0.7, [(no, 30, BROWN, True)], slide=s)
        text(x + 0.2, 2.75, cw - 0.4, 0.5, [(name, 18, INK, True)], slide=s)
        text(x + 0.2, 3.35, cw - 0.4, 1.8, [(desc, 12, GREY, False)], slide=s)
        text(x + 0.2, 5.45, cw - 0.4, 0.7, [("🛠 " + tool, 10, BROWN, False)], slide=s)

    # ---- P03 人群画像 ----
    s = new_slide()
    page_title(s, "品牌定位 —— 人群画像", "真实企划 deck 的入口：先定人群，再定产品结构")
    pw = 3.95
    for i, p in enumerate(personas[:3]):
        x = 0.55 + i * (pw + 0.15)
        card(x, 1.8, pw, 3.9, slide=s)
        text(x + 0.25, 2.1, pw - 0.5, 0.6, [(p["name"], 20, INK, True)], slide=s)
        text(x + 0.25, 2.85, pw - 0.5, 1.0, [(p.get("desc", ""), 14, GREY, False)], slide=s)
        text(x + 0.25, 4.4, pw - 0.5, 1.1, [("要什么：" + p.get("need", ""), 12, BROWN, False)], slide=s)

    # ---- P04 产品结构总表（企划盘聚合）----
    s = new_slide()
    page_title(s, "产品结构总表", f'{len(entries)} 个款式 ｜ 成分 / 规格 / 双价格')
    cols = [("品类", 2.4), ("成分", 3.3), ("规格", 2.6), ("会员价", 1.4), ("供货价", 1.4), ("对标", 2.1)]
    tx, ty, rh = 0.55, 1.8, 0.62
    box(tx, ty, 12.25, 0.5, fill=RGBColor(0xF0, 0xE9, 0xDC), line=LINE, slide=s)
    cx = tx
    for name, cwid in cols:
        text(cx + 0.12, ty + 0.09, cwid - 0.2, 0.35, [(name, 12, BROWN, True)], slide=s)
        cx += cwid
    ry = ty + 0.5
    for v in views:
        box(tx, ry, 12.25, rh, fill=WHITE, line=LINE, slide=s)
        row = [v["category"], v["composition"], v["spec"], v["msrp_str"], v["wholesale_str"], v["benchmark"]]
        cx = tx
        for val, (name, cwid) in zip(row, cols):
            text(cx + 0.12, ry + 0.16, cwid - 0.2, 0.4, [(str(val), 12, INK, name == "品类")], slide=s)
            cx += cwid
        ry += rh

    # ---- P05+ 逐款页 × N ----
    for idx, (e, v) in enumerate(zip(entries, views)):
        s = new_slide()
        sub = f'{v["color"] + " ｜ " if v["color"] else ""}{plan.get("direction","")} ｜ 对标 {v["benchmark"]} ｜ 逐款 {idx+1}/{len(entries)}'
        page_title(s, f'产品企划 · {v["category"]}', sub)
        layout = [t for t in LAYOUT_ORDER if t in e["selected"]]
        if layout:
            hero_key = "lifestyle" if "lifestyle" in e["selected"] else layout[0]
            fit_img(e["selected"][hero_key], 0.55, 1.7, 4.6, 4.4, slide=s)
            text(0.55, 6.25, 4.6, 0.4, [(TYPE_LABEL[hero_key], 11, BROWN, False)], slide=s)
            sec = [t for t in layout if t != hero_key][:1]
            if sec:
                fit_img(e["selected"][sec[0]], 5.3, 1.7, 1.9, 2.1, slide=s)
                text(5.3, 3.9, 1.9, 0.35, [(TYPE_LABEL[sec[0]].split(" ·")[0], 10, BROWN, False)], slide=s)
        px = 7.5
        card(px, 1.7, 5.3, 4.9, slide=s)
        rows = [("品名", v["category"]), ("成分", v["composition"]), ("规格", v["spec"]),
                ("细度", v["fineness"]), ("工艺", v["craft"]), ("功能", v["attributes"]),
                ("对标", v["benchmark"])]
        ry = 1.95
        for k, val in rows:
            if not val:
                continue
            text(px + 0.25, ry, 1.0, 0.35, [(k, 12, BROWN, True)], slide=s)
            text(px + 1.25, ry, 3.9, 0.35, [(val, 12, INK, False)], slide=s)
            ry += 0.42
        text(px + 0.25, ry + 0.1, 2.5, 0.6, [(v["msrp_str"], 20, INK, True),
                                             ("  会员价", 11, GREY, False)], slide=s)
        text(px + 2.9, ry + 0.1, 2.4, 0.6, [(v["wholesale_str"], 20, BROWN, True),
                                            ("  供货价", 11, GREY, False)], slide=s)
        dd = _design_dirs_str(v)
        if dd:
            text(px + 0.25, ry + 0.75, 4.9, 0.6, [("设计方向：" + dd, 11, GREY, False)], slide=s)

        # 演变对比（该款有记录时，紧跟其逐款页）
        ve = e.get("variation")
        if ve:
            s = new_slide()
            axis_label = AXIS_LABEL.get(ve.get("axis"), ve.get("axis", ""))
            page_title(s, f'演变对比 · {v["category"]}', f'{axis_label}：{ve.get("change", "")}')
            fit_img(ve["before"], 0.55, 1.7, 5.9, 5.0, slide=s)
            text(0.55, 6.8, 5.9, 0.4, [("畅销款（起点）", 12, BROWN, True)], slide=s)
            text(6.55, 3.9, 0.5, 0.6, [("→", 36, BROWN, True)], slide=s)
            fit_img(ve["after"], 7.15, 1.7, 5.9, 5.0, slide=s)
            text(7.15, 6.8, 5.9, 0.4, [("演变款（同族兄弟款，锁 DNA 只动一轴）", 12, BROWN, True)], slide=s)

    # ---- 品类矩阵：一页看全盘 ----
    s = new_slide()
    page_title(s, "品类矩阵 · 一页看全盘", f'{len(entries)} 个款式 × 首图 + 会员价')
    n = len(views)
    cols_n = min(n, 4)
    rows_n = (n + cols_n - 1) // cols_n
    gw, gh, ggap = min(2.9, 12.25 / cols_n - 0.15), 2.9, 0.15
    for i, (e, v) in enumerate(zip(entries, views)):
        r, c = divmod(i, cols_n)
        x = 0.55 + c * (gw + ggap)
        y = 1.75 + r * (gh + 0.75)
        first_img = next(iter(e["selected"].values()))
        card(x, y, gw, gh, slide=s)
        fit_img(first_img, x + 0.05, y + 0.05, gw - 0.1, gh - 0.1, slide=s)
        text(x, y + gh + 0.06, gw * 0.62, 0.35, [(v["category"], 12, INK, True)], slide=s)
        text(x + gw * 0.62, y + gh + 0.06, gw * 0.38, 0.35, [(v["msrp_str"], 13, BROWN, True)],
             align=PP_ALIGN.RIGHT, slide=s)

    # ---- AI 出图体系：行=款 × 列=图类型 ----
    s = new_slide()
    page_title(s, "AI 出图体系", "行=款 × 列=图类型 · 每类图独立锁定产品 DNA")
    ry = 1.8
    row_h = min(1.15, 5.4 / max(n, 1) - 0.15)
    for e, v in zip(entries, views):
        layout = [t for t in LAYOUT_ORDER if t in e["selected"]][:6]
        text(0.55, ry + row_h / 2 - 0.15, 1.6, 0.4, [(v["category"], 12, BROWN, True)], slide=s)
        iw = (12.25 - 1.8 - (max(len(layout), 1) - 1) * 0.1) / max(len(layout), 1)
        for j, t in enumerate(layout):
            x = 2.35 + j * (iw + 0.1)
            fit_img(e["selected"][t], x, ry, iw, row_h, slide=s)
        ry += row_h + 0.22

    # ---- 开发日历 ----
    s = new_slide()
    page_title(s, "FW26/27 开发日历", "企划 → 进店 8 节点")
    tw = (12.25 - 7 * 0.12) / 8
    for i, (d, name) in enumerate(CALENDAR):
        x = 0.55 + i * (tw + 0.12)
        card(x, 2.6, tw, 1.9, slide=s)
        text(x + 0.06, 2.85, tw - 0.12, 0.6, [(d, 11, BROWN, True)], slide=s)
        text(x + 0.06, 3.5, tw - 0.12, 0.9, [(name, 12, INK, True)], slide=s)
        if i < 7:
            text(x + tw + 0.0, 3.3, 0.14, 0.4, [("›", 14, BROWN, True)], slide=s)
    box(0.55, 5.6, 12.25, 0.85, fill=WHITE, line=LINE, slide=s)
    text(0.8, 5.78, 11.8, 0.5,
         [("⚡ AI 压缩：设计出图环节由「人天级」→「分钟级」——本 deck 全部概念图即 Agent 实跑产物。", 13, INK, False)], slide=s)

    # ---- 尾页 ----
    s = new_slide()
    text(0.9, 2.4, 11.5, 1.2, [("THANKS", 54, INK, True)], slide=s)
    text(0.9, 3.9, 11.5, 0.5, [("工具链：" + TOOLCHAIN, 13, GREY, False)], slide=s)
    text(0.9, 4.5, 11.5, 0.5, [(_timing_str(plan.get("timing")), 13, BROWN, False)], slide=s)
    text(0.9, 6.7, 11.5, 0.4, [(f"商品企划 Agent · 面向{plan.get('retailer', '')}的服装配饰 B端企划", 11, GREY, False)], slide=s)

    prs.save(out_path)
    return out_path
