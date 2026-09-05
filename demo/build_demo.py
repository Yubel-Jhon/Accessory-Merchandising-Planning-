#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""把 plan.json + 4 张 MVP 图 → 组装成「商品企划」demo 产物：
  1) demo/商品企划-山羊绒围巾.html  —— 自包含 live demo（base64 内嵌图）
  2) demo/商品企划-山羊绒围巾.pptx  —— 1 页推介 PPT

用法:  python build_demo.py
"""
import base64
import json
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEMO = os.path.dirname(os.path.abspath(__file__))

TYPE_LABEL = {
    "lifestyle": "氛围图 · 场景代入",
    "studio": "商拍图 · 质感定调",
    "white_bg": "白底图 · 产品身份证",
    "detail": "细节图 · 品质论证",
}


def b64(path):
    with open(path, "rb") as f:
        return "data:image/jpeg;base64," + base64.b64encode(f.read()).decode()


def load_plan():
    with open(os.path.join(DEMO, "plan.json"), encoding="utf-8") as f:
        return json.load(f)


def selection_cards(plan):
    labels = {
        "对标": "对标",
        "人群": "人群",
        "价格带": "价格带",
        "差异化": "差异化",
        "山姆锚点": "山姆锚点",
        "风险": "风险",
    }
    out = []
    for key, text in plan["selection_logic"].items():
        title = labels.get(key, key)
        out.append(
            f'<div class="card"><h3>{title}</h3><p>{text}</p></div>'
        )
    return "\n".join(out)


def build_html(plan, imgs):
    p = plan["product"]
    chips = [
        ("方向", plan["direction"]),
        ("品类", p["category"]),
        ("材质", p["material"]),
        ("工艺", p["craft"]),
        ("颜色", p["color"]),
        ("价格带", p["price_band"]),
        ("对标", p["benchmark"]),
        ("零售商", plan["retailer"]),
    ]
    chip_html = "".join(f'<span class="chip"><b>{k}</b>{v}</span>' for k, v in chips)

    grid = "".join(
        f'''<figure class="gitem {t}">
        <img src="{imgs[t]}" alt="{TYPE_LABEL[t]}"/>
        <figcaption>{TYPE_LABEL[t]}</figcaption>
        </figure>'''
        for t in plan["layout"]
    )

    html = """<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>山羊绒围巾 · 商品企划</title>
<style>
  :root{
    --ink:#2C2C2C; --brown:#8B7355; --camel:#C4A484;
    --paper:#F5F1EA; --card:#FFFFFF; --line:#E5DED2;
  }
  *{box-sizing:border-box;margin:0;padding:0}
  body{
    font-family:"Segoe UI","PingFang SC","Microsoft YaHei",sans-serif;
    background:var(--paper); color:var(--ink); line-height:1.65;
  }
  .wrap{max-width:1120px;margin:0 auto;padding:40px 28px 80px}
  header{text-align:center;padding:28px 0 8px}
  header .eyebrow{
    letter-spacing:.35em;font-size:12px;color:var(--brown);
    text-transform:uppercase;margin-bottom:10px;
  }
  header h1{font-family:Georgia,"Songti SC",serif;font-weight:500;font-size:40px;letter-spacing:.02em}
  header .sub{color:#7a6f5f;margin-top:10px;font-size:15px}
  .chips{display:flex;flex-wrap:wrap;gap:10px;justify-content:center;margin:26px 0 6px}
  .chip{background:var(--card);border:1px solid var(--line);border-radius:999px;
    padding:6px 14px;font-size:13px;color:#4a4238}
  .chip b{color:var(--brown);margin-right:6px;font-weight:600}
  .hero{margin:34px 0 26px;border-radius:6px;overflow:hidden;
    box-shadow:0 12px 40px rgba(44,44,44,.14)}
  .hero img{width:100%;display:block}
  h2.sec{font-family:Georgia,"Songti SC",serif;font-weight:500;font-size:22px;
    margin:46px 0 18px;display:flex;align-items:center;gap:12px}
  h2.sec::before{content:"";width:26px;height:2px;background:var(--brown);display:inline-block}
  .grid{display:grid;grid-template-columns:repeat(4,1fr);gap:14px}
  .gitem{background:var(--card);border:1px solid var(--line);border-radius:6px;
    overflow:hidden;display:flex;flex-direction:column}
  .gitem img{width:100%;aspect-ratio:1/1;object-fit:cover;display:block}
  .gitem figcaption{padding:10px 12px;font-size:13px;color:#5b5145;
    border-top:1px solid var(--line);background:#fbf9f5}
  .cards{display:grid;grid-template-columns:repeat(3,1fr);gap:14px}
  .card{background:var(--card);border:1px solid var(--line);border-radius:6px;padding:18px 20px}
  .card h3{font-size:14px;color:var(--brown);letter-spacing:.08em;margin-bottom:8px;
    font-weight:600;text-transform:uppercase}
  .card p{font-size:14px;color:#3f3930}
  footer{margin-top:60px;padding-top:20px;border-top:1px solid var(--line);
    font-size:12px;color:#a49a8b;text-align:center}
  @media(max-width:860px){.grid{grid-template-columns:repeat(2,1fr)}.cards{grid-template-columns:1fr}}
</style>
</head>
<body>
<div class="wrap">
  <header>
    <div class="eyebrow">Product Planning · Quiet Luxury</div>
    <h1>山羊绒围巾 · 商品企划</h1>
    <div class="sub">静奢 / 老钱风 ｜ 对标 Loro Piana ｜ 山姆会员锚点</div>
  </header>

  <div class="chips">__CHIPS__</div>

  <div class="hero"><img src="__IMG_LIFESTYLE__" alt="氛围图"></div>

  <h2 class="sec">图库编排 · 氛围 → 商拍 → 白底 → 细节</h2>
  <div class="grid">__GRID__</div>

  <h2 class="sec">选品逻辑</h2>
  <div class="cards">__CARDS__</div>

  <footer>由「商品企划生成 Agent」编排生成 ｜ 图类型：白底 / 商拍 / 氛围 / 细节 ｜ 生图：Seedream</footer>
</div>
</body>
</html>"""

    html = (html.replace("__CHIPS__", chip_html)
                .replace("__IMG_LIFESTYLE__", imgs["lifestyle"])
                .replace("__GRID__", grid)
                .replace("__CARDS__", selection_cards(plan)))
    return html


def build_pptx(plan, img_paths):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from PIL import Image

    INK = RGBColor(0x2C, 0x2C, 0x2C)
    BROWN = RGBColor(0x8B, 0x73, 0x55)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    def box(x, y, w, h, fill=None, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid() if fill else shp.fill.background()
        if fill:
            shp.fill.fore_color.rgb = fill
        if line:
            shp.line.color.rgb = line
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def text(x, y, w, h, runs, align=None, anchor=None):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        if anchor:
            tf.vertical_anchor = anchor
        for i, (txt, size, color, bold) in enumerate(runs):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.alignment = align
            r = para.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Microsoft YaHei"
        return tb

    def fit_img(path, cx, cy, cw, ch):
        """在 (cx,cy) 起始、宽高 (cw,ch) 的盒内等比例居中放图，返回实际宽高。"""
        with Image.open(path) as im:
            w, h = im.size
        ar = w / h
        box_ar = cw / ch
        if ar > box_ar:
            dw, dh = cw, cw / ar
        else:
            dh, dw = ch, ch * ar
        dx = cx + (cw - dw) / 2
        dy = cy + (ch - dh) / 2
        slide.shapes.add_picture(path, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
        return dw, dh

    # 背景
    box(0, 0, 13.333, 7.5, fill=RGBColor(0xF5, 0xF1, 0xEA))
    # 顶部色条
    box(0, 0, 13.333, 0.12, fill=BROWN)

    # 标题
    text(0.55, 0.28, 12.2, 0.9,
         [("山羊绒围巾 · 商品企划", 30, INK, True)],
         align=None)
    text(0.55, 0.95, 12.2, 0.4,
         [("静奢 / 老钱风   ·   对标 Loro Piana   ·   16.5um 小山羊绒   ·   山姆会员锚点", 13, BROWN, False)])

    # 左：hero 氛围图（16:9）
    fit_img(img_paths["lifestyle"], 0.55, 1.5, 7.4, 5.1)
    text(0.55, 6.62, 7.4, 0.4, [("氛围图 · 场景代入（卖给谁 / 什么场景）", 11, BROWN, False)])

    # 右：三张竖排（商拍 / 白底 / 细节）
    labels = [
        ("studio", "商拍图 · 质感定调"),
        ("white_bg", "白底图 · 产品身份证"),
        ("detail", "细节图 · 品质论证"),
    ]
    rx, rw = 8.25, 4.6
    th = 1.62
    gap = 0.12
    for i, (key, lab) in enumerate(labels):
        ty = 1.5 + i * (th + gap)
        box(rx, ty, rw, th, fill=RGBColor(0xFF, 0xFF, 0xFF), line=RGBColor(0xE5, 0xDE, 0xD2))
        fit_img(img_paths[key], rx + 0.06, ty + 0.06, th - 0.12, th - 0.12)
        text(rx + th + 0.1, ty + 0.18, rw - th - 0.2, 0.9,
             [(lab, 12, INK, True)],
             anchor=None)

    # 底部：选品逻辑一行
    box(0.55, 7.0, 12.25, 0.02, fill=RGBColor(0xE5, 0xDE, 0xD2))
    one = "选品逻辑：对标 Loro Piana 无 logo 羊绒围巾（专柜 ¥6000+），砍到 ¥299–499 山姆价格带；人群 = 山姆中产家庭「体面的性价比」；差异化 = 参数透明（16.5um / 12 针精纺）+ 无 logo 静奢。"
    text(0.55, 6.95, 12.25, 0.5, [(one, 10, RGBColor(0x5B, 0x51, 0x45), False)])

    out = os.path.join(DEMO, "商品企划-山羊绒围巾.pptx")
    prs.save(out)
    return out


def main():
    plan = load_plan()
    img_paths = {k: os.path.join(ROOT, v) for k, v in plan["images"].items()}
    for k, p in img_paths.items():
        if not os.path.exists(p):
            raise SystemExit(f"缺失图片: {p}")

    imgs = {k: b64(p) for k, p in img_paths.items()}

    html_out = os.path.join(DEMO, "商品企划-山羊绒围巾.html")
    with open(html_out, "w", encoding="utf-8") as f:
        f.write(build_html(plan, imgs))
    print("HTML ->", html_out)

    pptx_out = build_pptx(plan, img_paths)
    print("PPTX ->", pptx_out)


if __name__ == "__main__":
    main()
